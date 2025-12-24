import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.faiss import FAISS
from langchain.chains import RetrievalQA
from utils.config import llm, embeddings


# Directory to store persistent FAISS index
FAISS_INDEX_DIR = "faiss_index"


# Read and extract text from a PDF file
def _read_pdf(file):
    pdf = PdfReader(file)
    return "".join(page.extract_text() or "" for page in pdf.pages)


# Read and extract text from a Word document
def _read_docx(file):
    doc = Document(file)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


# Read and extract text from an Excel file
def _read_excel(file):
    dfs = pd.read_excel(file, sheet_name=None)
    text = ""

    for sheet_name, df in dfs.items():
        text += f"\nSheet: {sheet_name}\n"
        text += df.astype(str).to_string(index=False)

    return text


# Read and extract text from a PowerPoint file
def _read_pptx(file):
    prs = Presentation(file)
    text = ""

    for slide_no, slide in enumerate(prs.slides, start=1):
        text += f"\nSlide {slide_no}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


# Load existing FAISS index or create a new one
def _load_or_create_faiss(chunks):
    if os.path.exists(FAISS_INDEX_DIR):
        return FAISS.load_local(
            FAISS_INDEX_DIR,
            embeddings,
            allow_dangerous_deserialization=True
        )

    db = FAISS.from_texts(chunks, embeddings)
    db.save_local(FAISS_INDEX_DIR)
    return db


# Query content from multiple documents using semantic search
def query_documents(uploaded_files, query):

    # Validate uploaded documents
    if not uploaded_files:
        raise ValueError("No documents uploaded")

    # Ensure uniform list processing
    if not isinstance(uploaded_files, list):
        uploaded_files = [uploaded_files]

    all_text = ""

    # Extract text based on document type
    for file in uploaded_files:
        filename = file.name.lower()

        if filename.endswith(".pdf"):
            all_text += "\n" + _read_pdf(file)

        elif filename.endswith(".docx"):
            all_text += "\n" + _read_docx(file)

        elif filename.endswith(".xlsx"):
            all_text += "\n" + _read_excel(file)

        elif filename.endswith(".pptx"):
            all_text += "\n" + _read_pptx(file)

        else:
            raise ValueError(f"Unsupported file type: {file.name}")

    # Validate extracted text
    if not all_text.strip():
        raise ValueError("No readable text found in uploaded documents")

    # Split document text into overlapping chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=200
    )
    chunks = splitter.split_text(all_text)

    # Load or persist FAISS vector store
    db = _load_or_create_faiss(chunks)

    # Create retriever for semantic search
    retriever = db.as_retriever()

    # Initialize RetrievalQA chain
    qa_chain = RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever
    )

    # Execute query against indexed documents
    return qa_chain.invoke({"query": query})
